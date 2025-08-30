#!/usr/bin/env python3
# make_paper_slide.py — robust slide appender with theme/layout support
# - Precisely targets TITLE and BODY placeholders using PP_PLACEHOLDER
# - Works with arbitrary corporate themes
# - Atomic save with sanity checks

import argparse
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import PP_PLACEHOLDER as PH
import tempfile
import time
import sys


def _pick_placeholders(slide):
    """Return (title_shape, body_shape) using strict placeholder kinds, with fallbacks."""
    TITLE_TYPES = {
        getattr(PH, "TITLE", None),
        getattr(PH, "CENTER_TITLE", None),
        getattr(PH, "VERTICAL_TITLE", None),
    }
    BODY_TYPES = {
        getattr(PH, "BODY", None),
        getattr(PH, "OBJECT", None),
        getattr(PH, "VERTICAL_BODY", None),
        # Some themes misuse OBJECT as content; include it as body candidate.
    }
    EXCLUDE_AS_BODY = {
        getattr(PH, "TITLE", None),
        getattr(PH, "CENTER_TITLE", None),
        getattr(PH, "VERTICAL_TITLE", None),
        getattr(PH, "SUBTITLE", None),
        getattr(PH, "DATE", None),
        getattr(PH, "SLIDE_NUMBER", None),
        getattr(PH, "FOOTER", None),
        getattr(PH, "HEADER", None),
    }

    title_shape = None
    body_shape = None

    # First pass: strict by type
    for ph in slide.placeholders:
        if not getattr(ph, "is_placeholder", False):
            continue
        pht = ph.placeholder_format.type
        if pht in TITLE_TYPES and title_shape is None:
            title_shape = ph
            continue
        if pht in BODY_TYPES and body_shape is None:
            body_shape = ph
            continue

    # Fallbacks: choose by geometry if still missing
    if title_shape is None:
        text_phs = [ph for ph in slide.placeholders if hasattr(ph, "text_frame")]
        if text_phs:
            # Prefer highest (smallest 'top'); break ties by larger area
            title_shape = sorted(
                text_phs, key=lambda s: (s.top, -(s.width * s.height))
            )[0]

    if body_shape is None:
        candidates = []
        for ph in slide.placeholders:
            if ph is title_shape:
                continue
            if getattr(ph, "placeholder_format", None):
                if ph.placeholder_format.type in EXCLUDE_AS_BODY:
                    continue
            if hasattr(ph, "text_frame"):
                candidates.append(ph)
        if candidates:
            # Largest area that isn't the title/subtitle/date/footer/etc.
            body_shape = max(candidates, key=lambda s: (s.width * s.height))

    return title_shape, body_shape


def add_slide(prs, title, link, reason, usage, rel_label, app_label, layout_index):
    """Add a slide with a resolved title and structured body content."""
    # Obtain layout safely
    try:
        layout = prs.slide_layouts[layout_index]
    except Exception:
        print(f"⚠️ Layout index {layout_index} is invalid; falling back to 0", file=sys.stderr)
        layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(layout)

    # Identify placeholders robustly
    title_shape, body_shape = _pick_placeholders(slide)

    # --- Title: always in the title area (or injected top-left if absent)
    if title_shape is not None and hasattr(title_shape, "text_frame"):
        title_shape.text = title or "(title pending)"
        try:
            title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        except Exception:
            pass
    else:
        # No title placeholder in this layout; inject one at the top
        tbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(9.0), Inches(0.9))
        tf = tbox.text_frame
        tf.text = title or "(title pending)"
        tf.paragraphs[0].font.size = Pt(32)

    # --- Body: link + key bullets into the body placeholder (or injected box)
    if body_shape is not None and hasattr(body_shape, "text_frame"):
        body_tf = body_shape.text_frame
        body_tf.clear()
    else:
        # Create our own body area if theme layout doesn't provide one
        body_tf = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(6.7), Inches(5.0)).text_frame

    # Line 1: Link
    p0 = body_tf.paragraphs[0]
    p0.text = f"Link: {link}" if link else "Link: (n/a)"
    p0.font.size = Pt(14)

    # Subsequent lines: Relevance, Application (only if present)
    for label, val in [(rel_label, reason), (app_label, usage)]:
        if val:
            p = body_tf.add_paragraph()
            p.text = f"{label}: {val}"
            p.level = 1
            p.font.size = Pt(14)

    # Footer helper area (optional empty bullet for next action)
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(6.2), Inches(9.0), Inches(1.0))
    tf = tx.text_frame
    tf.text = "Next Action:"
    p = tf.add_paragraph()
    p.text = "• "
    p.level = 1


def save_atomic(prs: Presentation, dest: Path):
    """Save to a temp file and atomically replace the destination."""
    dest.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(prefix=dest.stem + "_", suffix=".pptx", dir=str(dest.parent), delete=False) as tmp:
        tmp_path = Path(tmp.name)
    try:
        prs.save(str(tmp_path))
        # Small sleep to avoid timestamp granularity issues on some filesystems
        time.sleep(0.05)
        os.replace(str(tmp_path), str(dest))
    except Exception as e:
        try:
            if tmp_path.exists():
                tmp_path.unlink()
        finally:
            raise e


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--title", required=True)
    ap.add_argument("--link", default="")
    ap.add_argument("--reason", default="")
    ap.add_argument("--usage", default="")
    ap.add_argument("--rel-label", default="Relevance")
    ap.add_argument("--app-label", default="Application")
    ap.add_argument("--outdir", default=str(Path.home() / "papers_slides"))
    ap.add_argument("--deck", default="")
    ap.add_argument("--theme", default=os.environ.get("PAPERFLOW_PPT_THEME", ""), help="Path to PPT theme .pptx")
    ap.add_argument(
        "--layout",
        type=int,
        default=int(os.environ.get("PAPERFLOW_PPT_LAYOUT", 1)),
        help="Layout index (0-based) inside the theme/deck",
    )
    args = ap.parse_args()

    # Load base presentation: existing deck > theme > blank
    prs = None
    if args.deck and Path(args.deck).exists():
        prs = Presentation(args.deck)
    elif args.theme and Path(args.theme).exists():
        prs = Presentation(args.theme)
    else:
        prs = Presentation()

    # Append slide
    before = len(prs.slides)
    add_slide(
        prs,
        title=args.title,
        link=args.link,
        reason=args.reason,
        usage=args.usage,
        rel_label=args.rel_label,
        app_label=args.app_label,
        layout_index=args.layout,
    )
    after = len(prs.slides)

    # Save: to deck if specified, otherwise create a single-slide file in outdir
    if args.deck:
        deck_path = Path(args.deck)
        save_atomic(prs, deck_path)
        # Reopen to verify
        try:
            prs_check = Presentation(str(deck_path))
            after_check = len(prs_check.slides)
        except Exception:
            after_check = -1
        print(f"Deck: {deck_path}")
        print(f"Slides: {before} -> {after} (reopen: {after_check})")
        try:
            print(f"Modified: {time.ctime(deck_path.stat().st_mtime)}")
        except Exception:
            pass
        print(str(deck_path))
    else:
        outdir = Path(args.outdir)
        outdir.mkdir(parents=True, exist_ok=True)
        # Safe file name based on title
        safe_name = args.title if args.title.strip() else "paper"
        safe_name = "".join(c if c.isalnum() or c in (" ", "_", "-") else "_" for c in safe_name)[:80]
        out = outdir / f"{safe_name}.pptx"
        save_atomic(prs, out)
        print("Slides: 0 -> 1")
        print(str(out))


if __name__ == "__main__":
    main()
